from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a PowerPoint presentation
prs = Presentation()

# Define the paths for each plot image on the new slide
slide1 = [
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/Effect Sizes/All/Outcome type_es_All.png',
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/MVPF/All/Outcome type_es_All.png'
]

# Set slide dimensions (increased height for a less stretched look)
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(8.5)  # Increased from 7.5 to 8.5 for extra height


# Add a blank slide layout
slide_layout = prs.slide_layouts[5]  # Use a blank layout
slide = prs.slides.add_slide(slide_layout)

# Add title with custom formatting for the new slide
title_shape = slide.shapes.title
title_shape.text = "Education PPPs show promising impacts on key outcomes, but with varying degrees of confidence as evidence remain limited for labor market outcomes"
#ducation PPPs Show Promising Impacts on Key Outcomes, with Varying Degrees of Confidence
# Adjust font style, size, color, and left alignment for the title
title_text_frame = title_shape.text_frame
title_paragraph = title_text_frame.paragraphs[0]
title_paragraph.font.size = Pt(24)
title_paragraph.font.name = 'Arial'
title_paragraph.font.color.rgb = RGBColor(26, 96, 110)  # Navy blue
title_paragraph.alignment = PP_ALIGN.LEFT

# Stretch the title bar almost to the end of the slide
title_shape.left = Inches(0.2)  # Small margin on the left
title_shape.width = Inches(13)  # Almost full width, leaving a small right margin
title_shape.top = Inches(0)  # Position at the top of the slide
title_shape.height = Inches(1)  # Set a reasonable height for the title bar

# Define positioning for the two images side by side
img_width = Inches(6)  # Width of each image
img_height = Inches(4.5)  # Height of each image
positions = [
    (Inches(0.5), Inches(2)),   # Left image
    (Inches(7), Inches(2))      # Right image
]

# Add the two images to the slide
for i, plot_path in enumerate(slide1):
    left, top = positions[i]
    slide.shapes.add_picture(plot_path, left, top, width=img_width, height=img_height)


# Define the paths for each plot image
plot_paths = [
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/Effect Sizes/All/Income classification_es_All.png',
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/Effect Sizes/Voucher/Income classification_es_Voucher.png',
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/Effect Sizes/Charter/Income classification_es_Charter.png',
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/MVPF/All/Income classification_es_All.png',
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/MVPF/Voucher/Income classification_es_Voucher.png',
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/MVPF/Charter/Income classification_es_Charter.png'
]


# Add a blank slide layout
slide_layout = prs.slide_layouts[5]  # Use a blank layout
slide = prs.slides.add_slide(slide_layout)

# Add title with custom formatting
title_shape = slide.shapes.title
title_shape.text = "Impacts are generally larger for developing countries, with disparities wider after considering costs"

# Adjust font style, size, color, and alignment
title_text_frame = title_shape.text_frame
title_paragraph = title_text_frame.paragraphs[0]
title_paragraph.font.size = Pt(24)
title_paragraph.font.name = 'Arial'
title_paragraph.font.color.rgb = RGBColor(26, 96, 110)  # Navy blue
title_paragraph.alignment = PP_ALIGN.LEFT

# Stretch the title bar almost to the end of the slide
title_shape.left = Inches(0.2)  # Small margin on the left
title_shape.width = Inches(13)  # Almost full width, leaving a small right margin
title_shape.top = Inches(0)  # Position at the top of the slide
title_shape.height = Inches(1)  # Set a reasonable height for the title bar

# Define positioning for the images in a 2x3 grid with lower row shifted down
img_width = Inches(4)  # Width of each image
img_height = Inches(3)  # Slightly increased height to adjust aspect ratio
positions = [
    (Inches(0.5), Inches(1.5)),  # Top-left
    (Inches(4.5), Inches(1.5)),  # Top-center
    (Inches(8.5), Inches(1.5)),  # Top-right
    (Inches(0.5), Inches(5)),    # Bottom-left, moved down
    (Inches(4.5), Inches(5)),    # Bottom-center, moved down
    (Inches(8.5), Inches(5))     # Bottom-right, moved down
]

# Add images to slide in adjusted 2x3 grid
for i, plot_path in enumerate(plot_paths):
    left, top = positions[i]
    slide.shapes.add_picture(plot_path, left, top, width=img_width, height=img_height)
    
plot_paths = [
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/Effect Sizes/All/Income classification_es_All.png',
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/Effect Sizes/Voucher/Income classification_es_Voucher.png',
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/Effect Sizes/Charter/Income classification_es_Charter.png',
    '/Users/angelosantos/Library/CloudStorage/OneDrive-GeorgeMasonUniversity-O365Production/School Choice/Outputs/MVPF/All/Income classification_es_All.png'
]

# Create a PowerPoint presentation
slide_layout = prs.slide_layouts[5]  # Use a blank layout

# Add a slide with the title
slide = prs.slides.add_slide(slide_layout)

# Add title with custom formatting
title_shape = slide.shapes.title
title_shape.text = "Impacts are generally larger for developing countries, with disparities wider after considering costs"

# Adjust font style, size, color, and alignment
title_text_frame = title_shape.text_frame
title_paragraph = title_text_frame.paragraphs[0]
title_paragraph.font.size = Pt(24)
title_paragraph.font.name = 'Arial'
title_paragraph.font.color.rgb = RGBColor(26, 96, 110)  # Navy blue
title_paragraph.alignment = PP_ALIGN.LEFT

# Stretch the title bar almost to the end of the slide
title_shape.left = Inches(0.2)  # Small margin on the left
title_shape.width = Inches(13)  # Almost full width, leaving a small right margin
title_shape.top = Inches(0)  # Position at the top of the slide
title_shape.height = Inches(1)  # Set a reasonable height for the title bar

# Define positioning for the images in a 2x3 grid with lower row shifted down

# Define positions for the images in a 2x2 grid
# Define positions for a 2x2 grid, adjusted to the right and lower
positions = [
    (Inches(1.5), Inches(1.5)),   # Top-left
    (Inches(6.47), Inches(1.5)),   # Top-right
    (Inches(1.5), Inches(5)),     # Bottom-left
    (Inches(6.47), Inches(5))      # Bottom-right
]


# Insert each image
for idx, plot_path in enumerate(plot_paths):
    slide.shapes.add_picture(plot_path, positions[idx][0], positions[idx][1], width=Inches(4.5), height=Inches(3.25))


# Save and overwrite the existing presentation
output_path = "/Users/angelosantos/Downloads/Impact_Slide_Presentation_upd.pptx"
prs.save(output_path)
